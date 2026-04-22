#!/usr/bin/env python3
from __future__ import annotations

import argparse
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.shapes.autoshape import Shape
from pptx.util import Inches, Pt


NAVY = RGBColor(12, 33, 64)
TEAL = RGBColor(0, 150, 136)
GRAY = RGBColor(120, 130, 140)
BG = RGBColor(248, 250, 252)
ASSET_DIR = Path(__file__).resolve().parents[1] / "assets"


def _set_run(run, *, size=24, bold=False, color=NAVY, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_title(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    # top bar
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.28))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), prs.slide_width - Inches(1.4), Inches(1.4))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run(r, size=44, bold=True, color=NAVY)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(12)
        r2 = p2.add_run()
        r2.text = subtitle
        _set_run(r2, size=20, color=GRAY)

    # visual placeholder box (will be replaced with generated image)
    ph = slide.shapes.add_picture(str(ASSET_DIR / "cover_climate_action.png"), Inches(8.3), Inches(1.1), Inches(4.3), Inches(3.2))


def add_bullets(prs: Presentation, title: str, bullets: list[str], note: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run(r, size=34, bold=True, color=NAVY)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), prs.slide_width - Inches(1.8), Inches(4.6))
    btf = box.text_frame
    btf.clear()
    btf.word_wrap = True
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(22)
        p.font.color.rgb = NAVY

    if note:
        n = slide.notes_slide.notes_text_frame
        n.text = note


def add_risk_identification(prs: Presentation, title: str, bullets: list[str], note: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run(r, size=34, bold=True, color=NAVY)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(6.0), Inches(4.6))
    btf = box.text_frame
    btf.clear()
    btf.word_wrap = True
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(22)
        p.font.color.rgb = NAVY

    img_path = str(ASSET_DIR / "diagram_risk_stack.png")
    slide.shapes.add_picture(img_path, Inches(7.5), Inches(2.2), width=Inches(5))

    if note:
        n = slide.notes_slide.notes_text_frame
        n.text = note


def add_profiles_comparison(prs: Presentation):
    """Key visual: same hazard, very different risk based on profile (exposure + vulnerability)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # title
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(0.9))
    tf = t.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "Same Hazard ≠ Same Risk"
    _set_run(r, size=36, bold=True, color=NAVY)

    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), prs.slide_width - Inches(1.6), Inches(0.6))
    st = subtitle.text_frame
    st.clear()
    r2 = st.paragraphs[0].add_run()
    r2.text = "Three profiles. Same heat hazard level. Very different risk."
    _set_run(r2, size=18, color=GRAY)

    # Three profile columns
    profiles = [
        ("Senior with respiratory issues", "High vulnerability", "Low adaptive capacity", "Risk: HIGH", RGBColor(200, 60, 60)),
        ("Young family in shaded home", "Medium vulnerability", "Good adaptive capacity", "Risk: MEDIUM", RGBColor(40, 140, 80)),
        ("Small business, poor drainage", "High exposure", "Low preparedness", "Risk: HIGH", RGBColor(200, 60, 60)),
    ]

    w = Inches(3.6)
    gap = Inches(0.5)
    start_x = (prs.slide_width - (3 * w + 2 * gap)) / 2

    # Image paths for profiles
    profile_images = [
        str(ASSET_DIR / "profile_senior_respiratory.png"),
        str(ASSET_DIR / "profile_family_shaded_home.png"),
        str(ASSET_DIR / "profile_business_poor_drainage.png"),
    ]

    for i, (name, v1, v2, risk, color) in enumerate(profiles):
        x = start_x + i * (w + gap)
        y = Inches(1.9)

        # Profile image
        img = slide.shapes.add_picture(profile_images[i], x + Inches(0.8), y + Inches(0.2), width=Inches(2.0))

        name_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.9), w - Inches(0.6), Inches(0.7))
        ntf = name_box.text_frame
        ntf.clear()
        nr = ntf.paragraphs[0].add_run()
        nr.text = name
        _set_run(nr, size=14, bold=True, color=NAVY)

        # Vulnerability factors
        v_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(2.7), w - Inches(0.6), Inches(1.8))
        vtf = v_box.text_frame
        vtf.clear()
        for line in [v1, v2]:
            p = vtf.add_paragraph()
            p.text = "• " + line
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(6)

        # Risk bar background
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.4), y + Inches(4.0), w - Inches(0.8), Inches(0.45))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = RGBColor(230, 235, 240)
        bar_bg.line.fill.background()

        # Risk level (different heights)
        heights = [0.85, 0.45, 0.82]  # relative heights
        bar_h = Inches(0.45 * heights[i])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.4), y + Inches(4.0) + Inches(0.45) - bar_h, w - Inches(0.8), bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Risk label
        risk_box = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(4.5), w - Inches(0.8), Inches(0.5))
        rtf = risk_box.text_frame
        rtf.clear()
        rr = rtf.paragraphs[0].add_run()
        rr.text = risk
        _set_run(rr, size=16, bold=True, color=color)

    # Bottom explanation
    expl = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), prs.slide_width - Inches(1.6), Inches(0.8))
    etf = expl.text_frame
    etf.clear()
    ep = etf.paragraphs[0]
    ep.text = "Our tool combines Hazard × Exposure × Vulnerability to give personalised risk — not just generic hazard maps."
    ep.font.size = Pt(15)
    ep.font.color.rgb = NAVY
    ep.font.bold = True
    ep.alignment = PP_ALIGN.CENTER


def add_actions_grid(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # title
    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = t.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "Actions to be taken (no‑regret first)"
    _set_run(r, size=32, bold=True, color=NAVY)

    # 2x2 grid + bottom strip
    cards = [
        ("Heat", ["Heat alert checklist", "Shade + ventilation", "Cool roofs / retrofits"]),
        ("Flood", ["Clear drains & gutters", "Backflow + barriers", "Drainage upgrades"]),
        ("Wildfire", ["Smoke plan + filters", "Defensible space", "Ember-proofing"]),
        ("Drought", ["Leak checks", "Efficiency & reuse", "Storage + demand mgmt"]),
    ]
    x0, y0 = Inches(0.8), Inches(1.55)
    w, h = Inches(4.45), Inches(2.3)
    gapx, gapy = Inches(0.5), Inches(0.45)

    for idx, (name, items) in enumerate(cards):
        col = idx % 2
        row = idx // 2
        x = x0 + col * (w + gapx)
        y = y0 + row * (h + gapy)
        rect = slide.shapes.add_shape(1, x, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(245, 247, 250)
        rect.line.color.rgb = RGBColor(220, 228, 236)

        hdr = slide.shapes.add_textbox(x + Inches(0.35), y + Inches(0.25), w - Inches(0.7), Inches(0.4))
        ht = hdr.text_frame
        ht.clear()
        rr = ht.paragraphs[0].add_run()
        rr.text = name
        _set_run(rr, size=22, bold=True, color=TEAL)

        body = slide.shapes.add_textbox(x + Inches(0.35), y + Inches(0.75), w - Inches(0.7), h - Inches(1.0))
        bt = body.text_frame
        bt.clear()
        for j, it in enumerate(items):
            p = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            p.text = it
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = NAVY

    # coastal strip
    strip = slide.shapes.add_shape(1, Inches(0.8), Inches(6.1), prs.slide_width - Inches(1.6), Inches(0.85))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(235, 245, 244)
    strip.line.color.rgb = RGBColor(200, 232, 228)
    st = slide.shapes.add_textbox(Inches(1.1), Inches(6.25), prs.slide_width - Inches(2.2), Inches(0.6))
    stf = st.text_frame
    stf.clear()
    p = stf.paragraphs[0]
    p.text = "Coastal storms: protect critical assets, plan for surge & access disruptions, coordinate early warnings"
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY


def build_deck(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 wide
    prs.slide_height = Inches(7.5)

    add_title(
        prs,
        "From Climate Signals to Action",
        "Identify local risks and prioritise practical actions (5‑minute pitch)",
    )

    add_bullets(
        prs,
        "Problem identification",
        [
            "Impacts are local, but information is fragmented and technical",
            "People need: “What should I do, here, first?”",
            "Without prioritisation: late reactions and misallocated budgets",
        ],
    )

    add_risk_identification(
        prs,
        "How we identify risk",
        [
            "Translate climate data into five hazard signals: heat, flood, wildfire, drought, coastal storm",
            "Combine with exposure and vulnerability (sensitivity + adaptive capacity)",
            "Explain the drivers so the response is targeted (not generic advice)",
        ],
    )

    add_profiles_comparison(prs)

    add_actions_grid(prs)

    add_bullets(
        prs,
        "Actions to take next (implementation)",
        [
            "Pilot with 1–3 municipalities / partners to validate local relevance",
            "Refine the measures library with domain experts and user feedback",
            "Produce a simple report: top risks, drivers, and the 3–5 highest-impact actions",
        ],
    )

    add_bullets(
        prs,
        "The ask",
        [
            "Pilot access (municipalities, civil protection, SME networks)",
            "Data + expertise for calibration and review",
            "Support to harden and deploy the tool at scale",
        ],
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main():
    ap = argparse.ArgumentParser(description="Generate a 5-minute innovation pitch deck (.pptx).")
    ap.add_argument(
        "--out",
        default=str(Path(__file__).resolve().parents[1] / "pitch_deck_final.pptx"),
        help="Output .pptx path",
    )
    args = ap.parse_args()
    build_deck(Path(args.out))
    print("Wrote", args.out)


if __name__ == "__main__":
    main()
